"""
RAGenie PowerPoint Presentation Generator
Generates a professional overview deck for GitHub / customer sharing.
Run: python generate_presentation.py
Output: RAGenie_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # very dark navy
CARD      = RGBColor(0x16, 0x2A, 0x3E)   # card background
CARD2     = RGBColor(0x1C, 0x34, 0x4F)   # slightly lighter card
CYAN      = RGBColor(0x00, 0xC8, 0xE0)   # primary accent
BLUE      = RGBColor(0x29, 0x8F, 0xFF)   # secondary blue
GREEN     = RGBColor(0x4A, 0xDE, 0x80)   # green accent
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)   # orange accent
PURPLE    = RGBColor(0xB3, 0x8A, 0xFF)   # purple accent
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # yellow accent
PINK      = RGBColor(0xFF, 0x5C, 0x8D)   # pink accent
TEAL      = RGBColor(0x00, 0xBF, 0xA5)   # teal accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xB0, 0xBE, 0xC5)
MGRAY     = RGBColor(0x55, 0x6B, 0x82)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor = None, line: RGBColor = None,
             line_w: int = 0, alpha: int = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()   # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w if line_w else 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color: RGBColor = WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  size=14, bold=False, color: RGBColor = WHITE,
                  align=PP_ALIGN.LEFT, spacing_after=4):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bld, sz, col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz if sz else size)
        run.font.bold = bld if bld is not None else bold
        run.font.color.rgb = col if col else color
    return txBox


def add_rounded_rect(slide, l, t, w, h, fill: RGBColor, line: RGBColor = None, line_w=1):
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def label_shape(slide, shape, text, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.5):
    """Draw a simple arrow connector using a line shape + arrowhead via XML."""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    # Add end arrowhead via XML
    ln = connector.line._ln
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'arrow')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')


def accent_bar(slide, color: RGBColor, width=0.07):
    add_rect(slide, 0, 0, width, 7.5, fill=color)


def section_header(slide, title, subtitle, accent: RGBColor):
    fill_bg(slide, BG)
    accent_bar(slide, accent)
    add_text(slide, title,   0.3, 0.25, 12.5, 1.0, size=36, bold=True,  color=WHITE,  align=PP_ALIGN.LEFT)
    add_text(slide, subtitle, 0.3, 1.2,  12.5, 0.6, size=18, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)


def divider_line(slide, y, color=MGRAY):
    add_rect(slide, 0.3, y, 12.73, 0.02, fill=color)

# ── Card helpers ───────────────────────────────────────────────────────────────

def feature_card(slide, l, t, w, h, icon, title, bullets, accent: RGBColor):
    """Draw a feature card with icon, title, and bullet list."""
    add_rounded_rect(slide, l, t, w, h, fill=CARD2, line=accent, line_w=1)
    # accent top bar
    add_rect(slide, l, t, w, 0.06, fill=accent)
    # icon + title
    add_text(slide, icon,  l+0.12, t+0.12, 0.5, 0.45, size=22, bold=False, color=accent)
    add_text(slide, title, l+0.55, t+0.12, w-0.65, 0.45, size=14, bold=True, color=WHITE)
    # bullets
    y_off = t + 0.62
    for b in bullets:
        add_text(slide, f"• {b}", l+0.15, y_off, w-0.25, 0.35, size=10.5, color=LGRAY)
        y_off += 0.38


def flow_box(slide, l, t, w, h, text, fill: RGBColor, text_size=10, text_color=WHITE, line: RGBColor = None):
    shape = add_rounded_rect(slide, l, t, w, h, fill=fill, line=line, line_w=1)
    label_shape(slide, shape, text, size=text_size, color=text_color)
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)

    # Gradient overlay strip (left accent)
    add_rect(s, 0, 0, 0.5, 7.5, fill=CYAN)
    add_rect(s, 0.5, 0, 0.5, 7.5, fill=BLUE)

    # Big glow circle (decorative)
    circ = s.shapes.add_shape(9, Inches(8.5), Inches(1.0), Inches(5.5), Inches(5.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x60)
    circ.line.fill.background()

    circ2 = s.shapes.add_shape(9, Inches(9.5), Inches(1.8), Inches(3.5), Inches(3.5))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x00, 0x60, 0x80)
    circ2.line.fill.background()

    # Logo text
    add_text(s, "RAGenie", 1.2, 1.8, 8, 1.6, size=72, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    add_text(s, "Your Local AI Assistant — Powered by RAG & LLMs",
             1.2, 3.5, 9, 0.7, size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, "100% Private  ·  On-Device  ·  Zero Cloud Dependency",
             1.2, 4.3, 9, 0.5, size=15, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)

    # Tag chips
    chips = [("🧠 Multi-LLM", BLUE), ("📄 RAG", TEAL), ("🔍 Web Search", GREEN),
             ("📊 Analytics", ORANGE), ("🔌 MCP", PURPLE), ("🔒 Auth", PINK)]
    x = 1.2
    for label, col in chips:
        add_rounded_rect(s, x, 5.1, 1.65, 0.5, fill=col)
        add_text(s, label, x+0.08, 5.18, 1.5, 0.35, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += 1.82

    add_text(s, "Open Source · MIT License", 1.2, 6.6, 6, 0.4, size=12, color=MGRAY)


def slide_what_is(prs):
    s = blank_slide(prs)
    section_header(s, "What is RAGenie?", "A self-hosted AI assistant that brings enterprise-grade capabilities to your local machine.", CYAN)
    divider_line(s, 1.9)

    desc = (
        "RAGenie combines Retrieval-Augmented Generation (RAG), local LLMs via Ollama, real-time "
        "internet search, data analytics, persistent memory, and task execution — all without sending "
        "a single byte to the cloud. Deploy it on your laptop or server and own your data completely."
    )
    add_text(s, desc, 0.4, 2.05, 12.5, 1.1, size=15, color=LGRAY, wrap=True)

    pillars = [
        ("🔒", "Privacy First",      "Everything runs locally.\nNo API keys, no telemetry.",         CYAN),
        ("⚡", "Real-time Streaming", "Token-by-token WebSocket\nstreaming for instant responses.",   BLUE),
        ("📚", "Document Knowledge",  "Upload PDFs, DOCX, Excel,\nimages, audio — ask anything.",     GREEN),
        ("🌐", "Live Web Search",     "DuckDuckGo integration with\n1-hour result cache.",             ORANGE),
        ("🧩", "MCP Ecosystem",       "Connects to 100s of MCP\nservers for tool use.",               PURPLE),
        ("📊", "Built-in Analytics",  "Statistical analysis, ML\nmodels, Plotly charts.",             PINK),
    ]
    x_positions = [0.3, 2.55, 4.8, 7.05, 9.3, 11.55]
    for (icon, title, body, col), x in zip(pillars, x_positions):
        add_rounded_rect(s, x, 3.3, 1.95, 3.7, fill=CARD2, line=col, line_w=1)
        add_rect(s, x, 3.3, 1.95, 0.07, fill=col)
        add_text(s, icon,  x+0.65, 3.45, 0.7, 0.55, size=26, align=PP_ALIGN.CENTER)
        add_text(s, title, x+0.1,  4.1,  1.75, 0.5, size=12, bold=True,  color=WHITE,  align=PP_ALIGN.CENTER)
        add_text(s, body,  x+0.1,  4.65, 1.75, 1.8, size=10, bold=False, color=LGRAY,  align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    s = blank_slide(prs)
    section_header(s, "System Architecture", "End-to-end flow from user request to streamed AI response.", BLUE)
    divider_line(s, 1.9)

    # ── Row 1: Frontend
    add_text(s, "FRONTEND  (React + TypeScript + Tailwind)", 0.3, 2.0, 12.7, 0.4, size=11, color=CYAN, bold=True)
    fe_boxes = [
        (0.3, "Chat UI\n+ File Upload"),
        (3.15, "Analytics\nDashboard"),
        (6.0, "Document\nSidebar"),
        (8.85, "MCP Servers\nManager"),
        (11.7, "News\nFeed"),
    ]
    for x, lbl in fe_boxes:
        flow_box(s, x, 2.45, 2.6, 0.75, lbl, BLUE, 9)

    # ── Arrows down
    for x in [1.6, 4.45, 7.3, 10.15]:
        add_arrow(s, x, 3.2, x, 3.65, CYAN, 1.2)

    # ── Row 2: Backend
    add_text(s, "BACKEND  (FastAPI + WebSocket)", 0.3, 3.65, 12.7, 0.35, size=11, color=ORANGE, bold=True)
    be_boxes = [
        (0.3,  "Chat\nOrchestrator"),
        (3.15, "Analytics\nEngine"),
        (6.0,  "Document\nManager"),
        (8.85, "MCP Client\nManager"),
        (11.7, "News\nAggregator"),
    ]
    for x, lbl in be_boxes:
        flow_box(s, x, 4.0, 2.6, 0.75, lbl, RGBColor(0x1A, 0x3A, 0x5C), text_size=9, line=ORANGE)

    # ── Arrows down
    for x in [1.6, 7.3]:
        add_arrow(s, x, 4.75, x, 5.2, CYAN, 1.2)
    add_arrow(s, 10.15, 4.75, 10.15, 5.2, PURPLE, 1.2)

    # ── Row 3: Core modules
    flow_box(s, 0.3,  5.2, 4.0, 0.7, "RAG System (BM25)  +  Multi-Model LLM Manager", CARD2, 9)
    flow_box(s, 4.5,  5.2, 4.0, 0.7, "Memory  ·  Search  ·  Tasks  ·  Learning", CARD2, 9)
    flow_box(s, 8.7,  5.2, 4.3, 0.7, "External MCP Servers  (100s of tools)", RGBColor(0x28, 0x12, 0x4A), 9)

    # ── Arrow down to Ollama
    add_arrow(s, 2.3, 5.9, 2.3, 6.35, GREEN, 1.2)

    # ── Row 4: Ollama
    flow_box(s, 0.3, 6.35, 8.5, 0.75,
             "Ollama  (Local LLMs)  ·  Reasoning Model  ·  Main Model  ·  Fallback Model", GREEN, 10)
    add_text(s, "100% On-Device", 9.0, 6.5, 2.5, 0.45, size=11, color=GREEN, bold=True)


def slide_rag_flow(prs):
    s = blank_slide(prs)
    section_header(s, "RAG Document Pipeline", "How documents are ingested, indexed, retrieved, and used in chat.", GREEN)
    divider_line(s, 1.9)

    # ── Left column: Ingestion
    add_text(s, "INGESTION", 0.3, 2.0, 4.0, 0.4, size=12, bold=True, color=GREEN)

    steps_ingest = [
        (GREEN,   "① Upload",        "PDF · DOCX · TXT · MD\nExcel · CSV · Image · Audio"),
        (TEAL,    "② Extract Text",  "pdfplumber · python-docx\nPillow OCR · SpeechRec"),
        (BLUE,    "③ Chunk",         "1000 char chunks\n200 char overlap"),
        (CYAN,    "④ SHA-256 Hash",  "Content-based dedup\nUnique chunk IDs"),
        (ORANGE,  "⑤ BM25 Index",   "Fast keyword indexing\npersisted to JSON"),
    ]
    y = 2.45
    for col, title, body in steps_ingest:
        add_rounded_rect(s, 0.3, y, 3.9, 0.82, fill=CARD2, line=col, line_w=1)
        add_text(s, title, 0.45, y+0.05, 3.6, 0.35, size=11, bold=True, color=col)
        add_text(s, body,  0.45, y+0.38, 3.6, 0.4,  size=9.5, color=LGRAY)
        if y < 6.0:
            add_arrow(s, 2.25, y+0.82, 2.25, y+0.97, GREEN, 1.2)
        y += 0.97

    # ── Center arrow
    add_rect(s, 4.35, 2.0, 0.05, 5.5, fill=MGRAY)
    add_text(s, "RETRIEVAL", 4.55, 2.0, 4.0, 0.4, size=12, bold=True, color=CYAN)

    # ── Right column: Query flow
    steps_query = [
        (CYAN,   "① User Query",    "Message arrives via\nWebSocket or REST"),
        (BLUE,   "② BM25 Search",  "Top-K chunks retrieved\nMedical term expansion"),
        (PURPLE, "③ Context Build", "Chunks assembled into\nstructured context block"),
        (ORANGE, "④ Prompt Build",  "System prompt + RAG\ncontext + user query"),
        (GREEN,  "⑤ LLM Response", "Streamed token-by-token\nback to frontend"),
    ]
    y = 2.45
    for col, title, body in steps_query:
        add_rounded_rect(s, 4.55, y, 3.9, 0.82, fill=CARD2, line=col, line_w=1)
        add_text(s, title, 4.7,  y+0.05, 3.6, 0.35, size=11, bold=True, color=col)
        add_text(s, body,  4.7,  y+0.38, 3.6, 0.4,  size=9.5, color=LGRAY)
        if y < 6.0:
            add_arrow(s, 6.5, y+0.82, 6.5, y+0.97, CYAN, 1.2)
        y += 0.97

    # ── Key design note
    add_rounded_rect(s, 8.65, 2.0, 4.5, 5.5, fill=CARD, line=MGRAY, line_w=1)
    add_text(s, "Design Decisions", 8.8, 2.1, 4.2, 0.45, size=13, bold=True, color=YELLOW)
    decisions = [
        ("BM25 vs Vectors", "Fast, dependency-light keyword\nsearch — no GPU embedding needed"),
        ("SHA-256 Dedup",   "Prevents re-ingesting the same\ncontent across uploads"),
        ("Overlapping Chunks", "200-char overlap preserves\ncontext at chunk boundaries"),
        ("30 MB limit",    "Upload cap enforced server-side\nwith rate limiting per IP"),
        ("BM25 persistence", "Index saved to JSON — survives\nserver restarts with no rebuild"),
    ]
    y = 2.6
    for title, body in decisions:
        add_text(s, f"▸ {title}", 8.8, y, 4.2, 0.3, size=10.5, bold=True,  color=CYAN)
        add_text(s, body,         8.8, y+0.3, 4.2, 0.45, size=9.5, color=LGRAY)
        y += 0.85


def slide_multi_model(prs):
    s = blank_slide(prs)
    section_header(s, "Multi-Model LLM Architecture",
                   "Three specialized models working in concert for quality, speed, and reliability.", ORANGE)
    divider_line(s, 1.9)

    # ── Three model cards
    models = [
        (PURPLE, "🧠", "Reasoning Model",  "deepseek-r1:1.5b",
         ["Chain-of-thought analysis", "Activates on complex queries", "Low temperature (0.3)", "1.1 GB — lightweight", "Auto-detected via keyword patterns"]),
        (CYAN,   "⚡", "Main Model",       "gemma4 / qwen2.5 / llama3.2",
         ["Primary response generation", "Receives reasoning context", "Standard temperature (0.7)", "Fully configurable via YAML", "Supports Ollama + HuggingFace"]),
        (GREEN,  "🛡️", "Fallback Model",   "qwen2.5:7b",
         ["Auto-activates on main failure", "Zero user intervention", "Ensures 100% uptime", "Independent config per role", "Timeout + error handling"]),
    ]
    x = 0.35
    for col, icon, title, model, bullets in models:
        add_rounded_rect(s, x, 2.05, 4.15, 5.2, fill=CARD2, line=col, line_w=1.5)
        add_rect(s, x, 2.05, 4.15, 0.08, fill=col)
        add_text(s, icon,  x+1.6,  2.18, 1.0, 0.55, size=28, align=PP_ALIGN.CENTER)
        add_text(s, title, x+0.15, 2.75, 3.85, 0.45, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, model, x+0.15, 3.2,  3.85, 0.4,  size=11, color=col, align=PP_ALIGN.CENTER)
        divider_line(s, 3.65)
        y = 3.75
        for b in bullets:
            add_text(s, f"✓  {b}", x+0.2, y, 3.75, 0.38, size=10.5, color=LGRAY)
            y += 0.42
        x += 4.42

    # ── Flow diagram at bottom
    add_text(s, "REQUEST ROUTING", 0.35, 7.0, 12.5, 0.35, size=10, bold=True, color=MGRAY)

    flow_box(s, 0.35, 7.35, 2.2, 0.5, "User Query", MGRAY, 10)
    add_arrow(s, 2.55, 7.6, 3.0, 7.6, CYAN)
    flow_box(s, 3.0, 7.35, 2.5, 0.5, "Auto-Detect\nReasoning?", BLUE, 9)
    # Yes branch
    add_text(s, "YES", 5.6, 7.42, 0.5, 0.28, size=9, color=GREEN, bold=True)
    add_arrow(s, 5.5, 7.6, 6.0, 7.6, PURPLE)
    flow_box(s, 6.0, 7.35, 2.1, 0.5, "Reasoning Model\n→ Analysis", PURPLE, 9)
    add_arrow(s, 8.1, 7.6, 8.6, 7.6, CYAN)
    # No branch text  
    add_text(s, "NO → directly to", 5.52, 7.75, 2.2, 0.3, size=8, color=MGRAY)
    flow_box(s, 8.6, 7.35, 2.1, 0.5, "Main Model\n→ Response", CYAN, 9)
    add_arrow(s, 10.7, 7.6, 11.2, 7.6, GREEN)
    flow_box(s, 11.2, 7.35, 1.9, 0.5, "Streamed\nto User", GREEN, 9)


def slide_capabilities(prs):
    s = blank_slide(prs)
    section_header(s, "Core Capabilities", "Everything RAGenie can do — at a glance.", TEAL)
    divider_line(s, 1.9)

    cards = [
        (CYAN,   "💬", "Chat & Memory",
         ["WebSocket token streaming", "Persistent conversation history", "Long-term memory (SQLite)", "Learning feedback loop", "Conversation pruning"]),
        (GREEN,  "📄", "Document RAG",
         ["PDF · DOCX · TXT · MD", "Excel · CSV · Images · Audio", "In-chat file upload (30 MB)", "BM25 keyword search", "Auto-deduplication"]),
        (ORANGE, "🌐", "Web Search",
         ["DuckDuckGo (no API key)", "1-hour result caching", "Real-time info injection", "Search result summarisation", "Integrated into RAG context"]),
        (PURPLE, "📊", "Data Analytics",
         ["CSV/Excel/JSON/TSV upload", "Statistical analysis (skew/kurtosis)", "Outlier detection (IQR/Z-score)", "Linear & RF regression/classif.", "Auto Plotly visualisations"]),
        (PINK,   "🔌", "MCP Integration",
         ["MCP server (SSE / stdio)", "MCP client manager UI", "Claude Desktop compatible", "stdio · SSE · HTTP transports", "Live tool rebuild on connect"]),
        (YELLOW, "🔒", "Security & Auth",
         ["JWT auth (access + refresh)", "Rate limiting per IP/WS", "Log redaction", "Security headers", "Audit log"]),
        (TEAL,   "📰", "News Aggregator",
         ["DuckDuckGo News (no key)", "Multi-language / region", "Auto LLM summarisation", "Ingest into RAG index", "3-day retention + cleanup"]),
        (BLUE,   "⚙️", "System",
         ["Single YAML config", "Env var overrides", "Hot-reload friendly", "Docker-ready structure", "MIT open-source license"]),
    ]

    positions = [
        (0.25, 2.1), (3.55, 2.1), (6.85, 2.1), (10.15, 2.1),
        (0.25, 5.05), (3.55, 5.05), (6.85, 5.05), (10.15, 5.05),
    ]
    for (col, icon, title, bullets), (x, y) in zip(cards, positions):
        add_rounded_rect(s, x, y, 3.05, 2.75, fill=CARD2, line=col, line_w=1)
        add_rect(s, x, y, 3.05, 0.06, fill=col)
        add_text(s, icon,  x+0.12, y+0.12, 0.5, 0.45, size=20, color=col)
        add_text(s, title, x+0.55, y+0.12, 2.4, 0.45, size=13, bold=True, color=WHITE)
        yb = y + 0.65
        for b in bullets:
            add_text(s, f"• {b}", x+0.15, yb, 2.8, 0.35, size=9.5, color=LGRAY)
            yb += 0.36


def slide_analytics(prs):
    s = blank_slide(prs)
    section_header(s, "Data Analytics Module", "Upload any data file — get instant statistical insights, ML models, and interactive charts.", ORANGE)
    divider_line(s, 1.9)

    # Pipeline flow
    pipeline = [
        (ORANGE, "Upload",      "CSV · Excel\nJSON · TSV · PDF"),
        (BLUE,   "Parse",       "pandas DataFrame\nall sheets + metadata"),
        (CYAN,   "Analyze",     "Basic + Advanced\nStats · Outliers"),
        (PURPLE, "ML Models",   "Regression\nClassification"),
        (GREEN,  "Visualize",   "Interactive\nPlotly Charts"),
        (YELLOW, "Predict",     "Trend Analysis\nFuture Values"),
    ]
    x = 0.3
    for col, title, body in pipeline:
        flow_box(s, x, 2.05, 2.0, 1.05, f"{title}\n{body}", col, 9)
        if x < 10.5:
            add_arrow(s, x+2.0, 2.58, x+2.35, 2.58, col, 1.5)
        x += 2.35

    # Stats detail box
    add_text(s, "Statistical Analysis", 0.3, 3.3, 6.0, 0.45, size=14, bold=True, color=ORANGE)
    stat_items = [
        "Mean · Median · Std · Min · Max · Quartiles",
        "Skewness · Kurtosis · IQR · Mode",
        "Correlation Matrix (Pearson)",
        "Outlier Detection: IQR Method + Z-Score Method",
        "Missing value detection & summary",
    ]
    y = 3.8
    for item in stat_items:
        add_text(s, f"✓  {item}", 0.3, y, 5.9, 0.35, size=11, color=LGRAY)
        y += 0.37

    # ML detail box
    add_text(s, "Machine Learning", 6.4, 3.3, 6.0, 0.45, size=14, bold=True, color=PURPLE)
    ml_items = [
        "Linear Regression — continuous target prediction",
        "Random Forest Regression — non-linear relationships",
        "Logistic Regression — binary classification",
        "Random Forest Classifier — multi-class classification",
        "Time Series — trend analysis & future forecasting",
    ]
    y = 3.8
    for item in ml_items:
        add_text(s, f"✓  {item}", 6.4, y, 6.5, 0.35, size=11, color=LGRAY)
        y += 0.37

    divider_line(s, 6.2)
    chart_types = ["Histogram", "Scatter", "Line", "Bar", "Box Plot", "Heatmap", "Pie Chart", "Auto-Viz"]
    add_text(s, "Chart Types:", 0.3, 6.3, 2.0, 0.4, size=11, bold=True, color=LGRAY)
    x = 2.5
    for ct in chart_types:
        add_rounded_rect(s, x, 6.28, 1.2, 0.4, fill=CARD2, line=ORANGE, line_w=1)
        add_text(s, ct, x+0.05, 6.32, 1.1, 0.32, size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        x += 1.35


def slide_mcp(prs):
    s = blank_slide(prs)
    section_header(s, "MCP Integration", "Model Context Protocol — exposing & consuming tools across the AI ecosystem.", PURPLE)
    divider_line(s, 1.9)

    # Left: MCP Server
    add_rounded_rect(s, 0.3, 2.1, 5.9, 5.1, fill=CARD, line=PURPLE, line_w=1)
    add_rect(s, 0.3, 2.1, 5.9, 0.08, fill=PURPLE)
    add_text(s, "🖥️  RAGenie as MCP SERVER", 0.5, 2.2, 5.5, 0.5, size=14, bold=True, color=PURPLE)
    server_items = [
        ("Transport", "SSE (port 8001)  or  stdio"),
        ("Protocol",  "mcp==1.17.0  (2025-06-18)"),
        ("Client",    "Claude Desktop, any MCP client"),
        ("Core tools", "search_documents · search_web\nask_ragenie · list_documents · execute_task"),
        ("News tools", "list/create/update/delete_news_keyword\nfetch_news_now · get_news_articles · suggest"),
        ("Mode",      "hybrid · mcp_server (headless)"),
    ]
    y = 2.85
    for k, v in server_items:
        add_text(s, k+":", 0.5,  y, 1.8, 0.4, size=10.5, bold=True, color=LGRAY)
        add_text(s, v,     2.4,  y, 3.6, 0.4, size=10.5, color=WHITE)
        y += 0.53

    # Right: MCP Client
    add_rounded_rect(s, 6.5, 2.1, 6.5, 5.1, fill=CARD, line=CYAN, line_w=1)
    add_rect(s, 6.5, 2.1, 6.5, 0.08, fill=CYAN)
    add_text(s, "🔌  RAGenie as MCP CLIENT", 6.7, 2.2, 6.1, 0.5, size=14, bold=True, color=CYAN)

    client_items = [
        ("Storage",    "SQLite DB for server configs"),
        ("Transports", "stdio · SSE · Streamable HTTP"),
        ("Tool naming", "server_name/tool_name (LLM-visible)"),
        ("UI",         "Manage via /mcp-servers Web UI"),
        ("Live rebuild", "Tools rebuild on connect/disconnect"),
        ("Deferred",   "Localhost servers auto-connect +10 s"),
    ]
    y = 2.85
    for k, v in client_items:
        add_text(s, k+":", 6.7,  y, 2.3,  0.4, size=10.5, bold=True, color=LGRAY)
        add_text(s, v,     9.1,  y, 3.75, 0.4, size=10.5, color=WHITE)
        y += 0.53

    # Ecosystem note
    add_text(s, "Compatible with the entire MCP ecosystem:", 6.7, 6.0, 6.2, 0.4, size=10.5, bold=True, color=LGRAY)
    tools = ["Filesystem", "GitHub", "Slack", "Databases", "Calendar", "Reminders", "Custom Tools"]
    x = 6.7
    y_t = 6.45
    for i, t in enumerate(tools):
        add_rounded_rect(s, x, y_t, 1.5, 0.38, fill=RGBColor(0x1A,0x2A,0x4A), line=CYAN, line_w=1)
        add_text(s, t, x+0.05, y_t+0.05, 1.4, 0.28, size=9, color=CYAN, align=PP_ALIGN.CENTER)
        x += 1.65
        if x > 12.5:
            x = 6.7
            y_t += 0.45


def slide_security(prs):
    s = blank_slide(prs)
    section_header(s, "Security & Authentication", "Production-grade security features built in from the ground up.", PINK)
    divider_line(s, 1.9)

    areas = [
        (PINK,   "🔑 JWT Authentication",
         ["Access token (30 min expiry)", "Refresh token (7 days)", "SQLite user store", "Enable/disable via config", "RAGENIE_SECRET_KEY env var"]),
        (ORANGE, "⏱️ Rate Limiting",
         ["60 req/min (general endpoints)", "10 uploads/hour per IP", "30 WS messages/min per client", "Configurable per endpoint type", "Protects from abuse"]),
        (YELLOW, "🛡️ Security Headers",
         ["X-Content-Type-Options", "X-Frame-Options: DENY", "Content-Security-Policy", "HSTS (Strict-Transport)", "Referrer-Policy"]),
        (CYAN,   "📝 Audit & Logging",
         ["Audit log at logs/audit.log", "Sensitive data redaction", "Structured JSON log format", "Configurable log level", "Rotation-ready"]),
        (PURPLE, "🔒 Request Controls",
         ["30 MB max upload size", "10 000 char WS message limit", "CORS origin whitelist", "File type validation", "Content-length enforcement"]),
        (GREEN,  "🏠 Privacy by Design",
         ["Zero cloud calls by default", "No telemetry or tracking", "All data stays on-device", "SQLite — no external DB", "Fully auditable open source"]),
    ]

    x, y = 0.3, 2.1
    for i, (col, title, bullets) in enumerate(areas):
        add_rounded_rect(s, x, y, 4.15, 2.55, fill=CARD2, line=col, line_w=1)
        add_rect(s, x, y, 4.15, 0.06, fill=col)
        add_text(s, title, x+0.15, y+0.1, 3.85, 0.45, size=12, bold=True, color=col)
        yb = y + 0.62
        for b in bullets:
            add_text(s, f"• {b}", x+0.15, yb, 3.85, 0.35, size=10, color=LGRAY)
            yb += 0.36
        x += 4.45
        if i == 2:
            x = 0.3
            y = 4.85


def slide_memory_learning(prs):
    s = blank_slide(prs)
    section_header(s, "Persistent Memory & Learning", "RAGenie remembers — and gets smarter with every interaction.", YELLOW)
    divider_line(s, 1.9)

    # Left: Memory
    add_rounded_rect(s, 0.3, 2.1, 6.0, 5.1, fill=CARD, line=YELLOW, line_w=1)
    add_rect(s, 0.3, 2.1, 6.0, 0.08, fill=YELLOW)
    add_text(s, "🧠  Persistent Memory", 0.5, 2.2, 5.6, 0.5, size=15, bold=True, color=YELLOW)
    mem_items = [
        "SQLite-backed long-term memory store",
        "Automatically stores important context from conversations",
        "Injects up to 8 memory items into every prompt",
        "2000-character context window per item",
        "Survives server restarts",
        "Proactive nudges based on stored user preferences",
        "Daily briefing at configured hour (default 9 AM)",
        "Quiet hours enforcement (10 PM – 8 AM default)",
    ]
    y = 2.85
    for m in mem_items:
        add_text(s, f"✓  {m}", 0.5, y, 5.6, 0.38, size=11, color=LGRAY)
        y += 0.44

    # Right: Learning
    add_rounded_rect(s, 6.6, 2.1, 6.0, 5.1, fill=CARD, line=GREEN, line_w=1)
    add_rect(s, 6.6, 2.1, 6.0, 0.08, fill=GREEN)
    add_text(s, "📈  Learning Feedback Loop", 6.8, 2.2, 5.6, 0.5, size=15, bold=True, color=GREEN)
    learn_items = [
        "Users give 👍/👎 on any AI response",
        "Positive feedback: score += 0.10",
        "Negative feedback: score -= 0.08",
        "Adaptation rate configurable in YAML",
        "Score influences future retrieval weight",
        "Per-session and aggregate feedback tracking",
        "REST API: POST /feedback/{message_id}",
        "Gradual system self-improvement over time",
    ]
    y = 2.85
    for m in learn_items:
        add_text(s, f"✓  {m}", 6.8, y, 5.6, 0.38, size=11, color=LGRAY)
        y += 0.44


def slide_news(prs):
    s = blank_slide(prs)
    section_header(s, "News Aggregator", "Stay informed — RAGenie fetches, summarises, and optionally indexes live news.", TEAL)
    divider_line(s, 1.9)

    # Pipeline
    news_flow = [
        (TEAL,   "Fetch",      "DuckDuckGo News\nNo API key required"),
        (BLUE,   "Filter",     "Region + language\nkeyword targeting"),
        (CYAN,   "Scrape",     "Full article text\n(up to 8000 chars)"),
        (ORANGE, "Summarise",  "LLM auto-summary\n(5-sentence default)"),
        (GREEN,  "Store",      "SQLite with 3-day\nretention policy"),
        (PURPLE, "RAG Inject", "Optional: ingest\ninto BM25 index"),
    ]
    x = 0.3
    for col, title, body in news_flow:
        flow_box(s, x, 2.1, 2.0, 1.1, f"{title}\n─────\n{body}", col, 9)
        if x < 10.5:
            add_arrow(s, x+2.0, 2.65, x+2.35, 2.65, col, 1.5)
        x += 2.35

    # Config grid
    add_text(s, "Configuration Options", 0.3, 3.45, 12.7, 0.45, size=14, bold=True, color=TEAL)
    config_rows = [
        ("keyword resolution",       "Plain-text term (e.g. 'IPL') — fuzzy match + LLM fallback to ID"),
        ("fetch interval",           "Every 60 minutes (configurable per keyword)"),
        ("articles per fetch",        "10 per topic (configurable 1–100)"),
        ("summarise on fetch",        "true — LLM summaries generated automatically"),
        ("ingest into RAG",           "false by default — enable to make articles searchable in chat"),
        ("retention_days",            "3 days — older articles auto-deleted on startup + every 6 h"),
        ("region",                    "wt-wt (worldwide) — supports 20+ language regions"),
    ]
    y = 3.95
    for k, v in config_rows:
        add_rounded_rect(s, 0.3, y, 12.7, 0.4, fill=CARD2)
        add_text(s, k,  0.5,  y+0.05, 3.5, 0.32, size=10.5, bold=True, color=TEAL)
        add_text(s, v,  4.1,  y+0.05, 8.8, 0.32, size=10.5, color=LGRAY)
        y += 0.45


def slide_tech_stack(prs):
    s = blank_slide(prs)
    section_header(s, "Technology Stack", "The open-source components that power RAGenie.", BLUE)
    divider_line(s, 1.9)

    stack = [
        ("Backend",       ORANGE, [
            ("FastAPI",       "Async REST + WebSocket server"),
            ("LangChain",     "LLM orchestration & tooling"),
            ("Ollama",        "Local LLM runtime (Metal/CUDA/CPU)"),
            ("HuggingFace",   "Alternative model provider"),
            ("BM25",         "Keyword-based document retrieval"),
            ("SQLite",        "Auth · Memory · News · Conversations"),
        ]),
        ("Frontend",      CYAN, [
            ("React 18",      "UI framework with hooks"),
            ("TypeScript",    "Type-safe frontend"),
            ("Vite",          "Dev server + production bundler"),
            ("Tailwind CSS",  "Utility-first styling"),
            ("Plotly.js",     "Interactive data visualisations"),
            ("Lucide Icons",  "Consistent icon set"),
        ]),
        ("Data & ML",     GREEN, [
            ("pandas",        "Data loading & manipulation"),
            ("scikit-learn",  "ML models (regression/classification)"),
            ("pdfplumber",    "PDF text extraction"),
            ("python-docx",   "DOCX/DOC processing"),
            ("Pillow + pytesseract", "Image metadata + OCR"),
            ("mutagen + SpeechRec",  "Audio metadata + transcription"),
        ]),
        ("Integration",   PURPLE, [
            ("MCP SDK",       "mcp==1.17.0 (2025-06-18 protocol)"),
            ("DuckDuckGo",    "Web search + news (no API key)"),
            ("python-jose",   "JWT token generation/validation"),
            ("passlib",       "Bcrypt password hashing"),
            ("slowapi",       "Rate limiting middleware"),
            ("aiofiles",      "Async file I/O"),
        ]),
    ]

    x = 0.3
    for category, col, items in stack:
        add_rounded_rect(s, x, 2.1, 3.1, 5.1, fill=CARD2, line=col, line_w=1)
        add_rect(s, x, 2.1, 3.1, 0.08, fill=col)
        add_text(s, category, x+0.15, 2.2, 2.8, 0.45, size=14, bold=True, color=col)
        y = 2.75
        for lib, desc in items:
            add_text(s, lib,  x+0.15, y,      2.8, 0.32, size=10.5, bold=True, color=WHITE)
            add_text(s, desc, x+0.15, y+0.28, 2.8, 0.32, size=9.5,  color=LGRAY)
            y += 0.65
        x += 3.3


def slide_api(prs):
    s = blank_slide(prs)
    section_header(s, "API Reference", "REST endpoints + WebSocket — integrable into any application.", CYAN)
    divider_line(s, 1.9)

    sections = [
        ("Chat & Documents", CYAN, [
            ("GET",    "/health",             "Health check"),
            ("POST",   "/chat",               "Send message, get AI response"),
            ("WS",     "/ws/{client_id}",     "Streaming chat (token-by-token)"),
            ("GET",    "/documents",           "List ingested documents"),
            ("POST",   "/upload",             "Upload + ingest document"),
            ("POST",   "/chat-upload",        "In-chat file upload (ingest + preview)"),
            ("DELETE", "/documents/{id}",     "Remove document from index"),
        ]),
        ("News Aggregator", TEAL, [
            ("GET",    "/api/keywords",                "List tracked keywords"),
            ("POST",   "/api/keywords",                "Create / track a keyword"),
            ("PATCH",  "/api/keywords/{id}",           "Update keyword settings"),
            ("DELETE", "/api/keywords/{id}",           "Delete keyword"),
            ("POST",   "/api/keywords/{id}/fetch-now", "Trigger immediate fetch"),
            ("GET",    "/api/news",                    "List articles (paginated)"),
            ("POST",   "/api/keywords/suggest",        "LLM keyword suggestion"),
        ]),
        ("MCP Client", PURPLE, [
            ("GET",    "/mcp-servers",              "List registered MCP servers"),
            ("POST",   "/mcp-servers",              "Register new MCP server"),
            ("GET",    "/mcp-servers/{id}",         "Get server details"),
            ("PUT",    "/mcp-servers/{id}",         "Update server config"),
            ("DELETE", "/mcp-servers/{id}",         "Remove MCP server"),
            ("POST",   "/mcp-servers/{id}/connect", "Connect to server"),
            ("POST",   "/mcp-servers/chat",         "Agent chat via all MCP tools"),
        ]),
    ]

    METHOD_COLORS = {
        "GET":    RGBColor(0x4A, 0xDE, 0x80),
        "POST":   BLUE,
        "DELETE": RGBColor(0xFF, 0x5C, 0x5C),
        "PUT":    YELLOW,
        "PATCH":  ORANGE,
        "WS":     CYAN,
    }

    x = 0.3
    for sec_title, col, rows in sections:
        add_text(s, sec_title, x, 2.0, 4.2, 0.4, size=12, bold=True, color=col)
        y = 2.45
        for method, path, desc in rows:
            add_rounded_rect(s, x, y, 4.2, 0.58, fill=CARD2)
            mc = METHOD_COLORS.get(method, LGRAY)
            add_rounded_rect(s, x+0.07, y+0.1, 0.65, 0.36, fill=mc)
            add_text(s, method, x+0.07, y+0.12, 0.65, 0.3, size=8.5, bold=True, color=BG, align=PP_ALIGN.CENTER)
            add_text(s, path,   x+0.77, y+0.05, 3.3, 0.3, size=9.5, bold=True, color=WHITE)
            add_text(s, desc,   x+0.77, y+0.3,  3.3, 0.25, size=9,  color=LGRAY)
            y += 0.65
        x += 4.45


def slide_getting_started(prs):
    s = blank_slide(prs)
    section_header(s, "Getting Started", "Up and running in under 5 minutes.", GREEN)
    divider_line(s, 1.9)

    steps = [
        ("1", GREEN,  "Clone & Install",
         "git clone https://github.com/ratulsarkar-iam/RAGenie.git\ncd RAGenie\npython3 -m venv venv && source venv/bin/activate\npip install -r requirements.txt"),
        ("2", ORANGE, "Install Ollama + Pull a Model",
         "# Install Ollama from https://ollama.com\nollama pull llama3.2          # 2 GB — quick start\nollama pull qwen2.5:7b        # 4.7 GB — recommended\nollama pull deepseek-r1:1.5b  # 1.1 GB — reasoning"),
        ("3", CYAN,   "Create Directories & Start Backend",
         "mkdir -p data/documents data/index logs\npython run_server.py\n# API live at http://localhost:8000"),
        ("4", PURPLE, "Start Frontend",
         "cd frontend\nnpm install\nnpm run dev\n# UI live at http://localhost:5173"),
    ]

    x = 0.3
    for num, col, title, code in steps:
        add_rounded_rect(s, x, 2.1, 3.1, 5.05, fill=CARD2, line=col, line_w=1)
        add_rect(s, x, 2.1, 3.1, 0.07, fill=col)
        # Step circle
        circ = s.shapes.add_shape(9, Inches(x+1.2), Inches(2.25), Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
        label_shape(s, circ, num, size=20, bold=True, color=BG)
        add_text(s, title, x+0.1, 3.1, 2.9, 0.45, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Code block
        add_rounded_rect(s, x+0.1, 3.6, 2.9, 1.4, fill=RGBColor(0x07, 0x11, 0x1C))
        add_text(s, code, x+0.18, 3.65, 2.75, 1.32, size=8.5, color=GREEN, wrap=True)
        x += 3.3

    # One-command note
    add_rounded_rect(s, 0.3, 7.25, 12.7, 0.55, fill=RGBColor(0x0A,0x28,0x1A), line=GREEN, line_w=1)
    add_text(s, "⚡  One-command start: ", 0.5, 7.33, 3.0, 0.38, size=12, bold=True, color=GREEN)
    add_text(s, "./start.sh  — launches both backend and frontend together", 3.5, 7.33, 9.2, 0.38, size=12, color=WHITE)


def slide_requirements(prs):
    s = blank_slide(prs)
    section_header(s, "System Requirements", "Minimum and recommended specs to run RAGenie comfortably.", LGRAY)
    divider_line(s, 1.9)

    reqs = [
        ("Python",  "3.9+",   "3.11+",  "Core backend runtime"),
        ("Node.js", "16+",    "18+",     "Frontend build toolchain"),
        ("RAM",     "8 GB",   "16 GB",   "Needed for running LLMs locally"),
        ("Disk",    "5 GB",   "20 GB",   "Space for models + documents"),
        ("Ollama",  "Latest", "Latest",  "Local LLM inference engine"),
        ("GPU/NPU", "CPU only","Apple M1+ / NVIDIA CUDA","Hardware acceleration (optional)"),
    ]

    # Header row
    add_rect(s, 0.3, 2.05, 12.7, 0.52, fill=CARD2)
    for col_x, label, w in [(0.5,  "Component", 2.5), (3.2, "Minimum", 2.5),
                             (5.9,  "Recommended", 3.0), (9.1, "Notes", 3.8)]:
        add_text(s, label, col_x, 2.1, w, 0.4, size=12, bold=True, color=LGRAY)

    y = 2.6
    for i, (comp, mn, rec, note) in enumerate(reqs):
        bg = CARD if i % 2 == 0 else CARD2
        add_rect(s, 0.3, y, 12.7, 0.55, fill=bg)
        add_text(s, comp, 0.5, y+0.1, 2.5, 0.38, size=12, bold=True, color=WHITE)
        add_text(s, mn,   3.2, y+0.1, 2.5, 0.38, size=12, color=ORANGE)
        add_text(s, rec,  5.9, y+0.1, 3.0, 0.38, size=12, color=GREEN)
        add_text(s, note, 9.1, y+0.1, 3.8, 0.38, size=11, color=LGRAY)
        y += 0.58

    # OS support
    add_text(s, "Supported Platforms", 0.3, 6.1, 4.0, 0.45, size=13, bold=True, color=LGRAY)
    platforms = [("🍎 macOS", "M1/M2/M3 — MPS acceleration", GREEN),
                 ("🐧 Linux", "CUDA/CPU — full support",      BLUE),
                 ("🪟 Windows", "CPU — Docker recommended",   ORANGE)]
    x = 0.3
    for icon_title, detail, col in platforms:
        add_rounded_rect(s, x, 6.6, 4.1, 0.7, fill=CARD2, line=col, line_w=1)
        add_text(s, icon_title, x+0.15, 6.65, 3.8, 0.35, size=12, bold=True, color=col)
        add_text(s, detail,     x+0.15, 6.98, 3.8, 0.3,  size=10, color=LGRAY)
        x += 4.35


def slide_differentiators(prs):
    s = blank_slide(prs)
    section_header(s, "Why RAGenie?", "Key differentiators over cloud-based alternatives.", CYAN)
    divider_line(s, 1.9)

    rows = [
        ("Feature",              "RAGenie",                 "ChatGPT / Gemini",    "Self-hosted LLM solutions"),
        ("Privacy",              "✅ 100% on-device",       "❌ Data sent to cloud","⚠️  Varies"),
        ("Cost",                 "✅ Free forever",         "❌ Monthly subscription","✅ Free"),
        ("Document RAG",         "✅ BM25 + 8 file types", "⚠️  Limited uploads",  "⚠️  Manual setup"),
        ("Multi-Model",          "✅ 3-role orchestration", "❌ Single model",      "⚠️  Manual"),
        ("Analytics",            "✅ Built-in + ML models", "❌ Not included",      "❌ Not included"),
        ("MCP Integration",      "✅ Server + Client",      "⚠️  Via ChatGPT plugins","❌ Not standard"),
        ("News Aggregator",      "✅ Built-in",             "❌ Not included",      "❌ Not included"),
        ("Persistent Memory",    "✅ SQLite-backed",        "⚠️  Session only",     "⚠️  Varies"),
        ("Auth & Rate Limiting", "✅ JWT + slowapi",        "✅ Managed",           "⚠️  Manual"),
        ("Open Source",          "✅ MIT License",          "❌ Proprietary",       "✅ Varies"),
    ]

    col_widths = [3.5, 3.2, 3.2, 3.2]
    col_x      = [0.3, 3.9, 7.2, 10.5]
    header_colors = [MGRAY, CYAN, LGRAY, LGRAY]

    # Header
    add_rect(s, 0.3, 2.05, 12.7, 0.52, fill=CARD2)
    for cx, w, label, col in zip(col_x, col_widths, rows[0], header_colors):
        add_text(s, label, cx+0.1, 2.1, w, 0.4, size=12, bold=True, color=col)

    y = 2.6
    for i, row in enumerate(rows[1:]):
        bg = CARD if i % 2 == 0 else CARD2
        add_rect(s, 0.3, y, 12.7, 0.49, fill=bg)
        for cx, w, cell, fc in zip(col_x, col_widths, row,
                                    [WHITE, CYAN, LGRAY, LGRAY]):
            add_text(s, cell, cx+0.1, y+0.08, w, 0.35, size=10.5, color=fc)
        y += 0.5


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)

    # Decorative circles
    for cx, cy, sz, col in [(10.5, 1.0, 4.5, RGBColor(0x00,0x30,0x50)),
                             (11.5, 2.5, 2.5, RGBColor(0x00,0x50,0x70))]:
        c = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    add_text(s, "RAGenie", 0.8, 1.5, 10, 1.4, size=72, bold=True, color=CYAN)
    add_text(s, "The AI assistant that stays on YOUR machine.", 0.8, 3.1, 10, 0.7, size=22, color=WHITE)
    add_text(s, "100% Private  ·  Fully Open Source  ·  Zero Cloud  ·  MIT License",
             0.8, 3.9, 11, 0.55, size=16, color=LGRAY)

    links = [
        ("🔗 GitHub",     "https://github.com/ratulsarkar-iam/RAGenie", CYAN),
        ("📄 Docs",       "README.md included in repository",            BLUE),
        ("⚡ Quick Start", "./start.sh  to launch everything",          GREEN),
        ("📧 License",    "MIT — free for commercial & personal use",    ORANGE),
    ]
    x = 0.8
    for label, detail, col in links:
        add_rounded_rect(s, x, 5.0, 2.9, 0.95, fill=CARD2, line=col, line_w=1)
        add_text(s, label,  x+0.15, 5.1,  2.6, 0.38, size=13, bold=True, color=col)
        add_text(s, detail, x+0.15, 5.5,  2.6, 0.38, size=9.5, color=LGRAY)
        x += 3.15

    add_text(s, "Built with ❤️  using FastAPI · LangChain · React · Ollama · MCP SDK 1.17.0",
             0.8, 6.8, 11.5, 0.45, size=12, color=MGRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# BLUEPRINT / WORKFLOW SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def _wf_box(s, l, t, w, h, text, fill, size=10, text_col=WHITE, line=None):
    """Process step — rounded rectangle."""
    shape = add_rounded_rect(s, l, t, w, h, fill=fill, line=line, line_w=1)
    label_shape(s, shape, text, size=size, color=text_col)
    return shape


def _wf_diamond(s, l, t, w, h, text, fill, size=9.5, text_col=WHITE):
    """Decision node — diamond shape."""
    shape = s.shapes.add_shape(4, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    label_shape(s, shape, text, size=size, color=text_col, align=PP_ALIGN.CENTER)
    return shape


def _wf_oval(s, l, t, w, h, text, fill, size=10, text_col=WHITE):
    """Start/End terminal — oval."""
    shape = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    label_shape(s, shape, text, size=size, bold=True, color=text_col)
    return shape


def _wf_label(s, text, l, t, col=LGRAY, size=8.5):
    add_text(s, text, l, t, 0.9, 0.3, size=size, color=col, align=PP_ALIGN.CENTER)


def _wf_arrow(s, x1, y1, x2, y2, col=CYAN, w=1.5):
    add_arrow(s, x1, y1, x2, y2, col, w)


# ── Blueprint 1 : Chat Request Lifecycle ─────────────────────────────────────

def slide_wf_chat(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, CYAN, 0.07)

    add_text(s, "WORKFLOW  ①", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=CYAN)
    add_text(s, "Chat Request Lifecycle", 0.3, 0.45, 12, 0.85, size=32, bold=True, color=WHITE)
    add_text(s, "Complete end-to-end flow from user input to streamed AI response",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # ── Left lane: Receive & Classify (x≈0.3 → 3.5)
    # ── Mid lane: Retrieve & Build (x≈3.5 → 7.5)
    # ── Right lane: LLM & Stream (x≈7.5 → 13.0)

    lane_ys = [0.07, 0.07, 0.07]
    lane_cols = [RGBColor(0x06,0x18,0x2A), RGBColor(0x06,0x1E,0x32), RGBColor(0x06,0x1A,0x2C)]
    lane_labels = ["RECEIVE & AUTH", "RETRIEVE & BUILD", "GENERATE & STREAM"]
    lane_label_cols = [CYAN, GREEN, ORANGE]
    lane_xs = [0.22, 4.45, 8.68]
    lane_widths = [4.0, 4.0, 4.45]
    for lx, lw, lbl, col, lc in zip(lane_xs, lane_widths, lane_labels, lane_cols, lane_label_cols):
        add_rect(s, lx, 1.85, lw, 5.55, fill=col)
        add_text(s, lbl, lx+0.05, 1.88, lw-0.1, 0.32, size=9, bold=True,
                 color=lc, align=PP_ALIGN.CENTER)

    # ══ Lane 1: Receive & Auth ══════════════════════════════
    CX = 2.22   # center x of lane 1
    _wf_oval(s,   CX-0.8, 2.25, 1.6, 0.48, "User Sends Message", CYAN, 9)
    _wf_arrow(s,  CX, 2.73, CX, 3.05, CYAN)
    _wf_box(s,   CX-1.0, 3.05, 2.0, 0.5, "WebSocket / REST\nHandler", BLUE, 9)
    _wf_arrow(s,  CX, 3.55, CX, 3.85, CYAN)
    _wf_diamond(s, CX-0.85, 3.85, 1.7, 0.6, "Auth\nEnabled?", RGBColor(0x1A,0x3A,0x5C), 9)
    # Yes → validate JWT
    _wf_label(s, "YES", CX+0.9, 4.05, GREEN)
    _wf_arrow(s, CX+0.85, 4.15, CX+1.45, 4.15, GREEN)
    _wf_box(s, CX+1.45, 3.9, 1.8, 0.5, "Validate JWT\nToken", RGBColor(0x0A,0x28,0x18), 9, GREEN)
    # No → continue
    _wf_label(s, "NO", CX-1.2, 4.05, ORANGE)
    _wf_arrow(s, CX, 4.45, CX, 4.75, CYAN)
    _wf_box(s, CX-1.0, 4.75, 2.0, 0.5, "Rate Limit\nCheck", RGBColor(0x1A,0x1A,0x3A), 9)
    _wf_arrow(s, CX, 5.25, CX, 5.55, CYAN)
    _wf_diamond(s, CX-0.85, 5.55, 1.7, 0.6, "Reasoning\nMode?", RGBColor(0x2A,0x10,0x40), 9)
    _wf_label(s, "AUTO-DETECT", CX+0.88, 5.75, PURPLE, 7.5)

    # Cross-lane arrow to Lane 2
    _wf_arrow(s, CX+0.85, 5.85, 4.45, 5.85, CYAN, 1.2)

    # ══ Lane 2: Retrieve & Build ════════════════════════════
    MX = 6.45   # center x of lane 2
    _wf_box(s, MX-1.1, 2.25, 2.2, 0.5, "Parse Message\n+ File Attachments", RGBColor(0x0A,0x28,0x18), 9, GREEN)
    _wf_arrow(s, MX, 2.75, MX, 3.05, GREEN)
    _wf_box(s, MX-1.1, 3.05, 2.2, 0.5, "BM25 Search\nDocument Index", RGBColor(0x0A,0x28,0x28), 9, TEAL)
    _wf_arrow(s, MX, 3.55, MX, 3.85, GREEN)
    _wf_box(s, MX-1.1, 3.85, 2.2, 0.5, "Retrieve Top-K\nDocument Chunks", CARD2, 9)
    _wf_arrow(s, MX, 4.35, MX, 4.65, GREEN)
    _wf_box(s, MX-1.1, 4.65, 2.2, 0.5, "DuckDuckGo\nWeb Search", RGBColor(0x18,0x28,0x0A), 9, GREEN)
    _wf_arrow(s, MX, 5.15, MX, 5.45, GREEN)
    _wf_box(s, MX-1.1, 5.45, 2.2, 0.5, "Inject Memory\nContext (top-8)", RGBColor(0x28,0x20,0x08), 9, YELLOW)
    _wf_arrow(s, MX, 5.95, MX, 6.25, GREEN)
    _wf_box(s, MX-1.1, 6.25, 2.2, 0.5, "Build Final Prompt\n(RAG + Memory + Query)", CARD2, 9)

    # Cross-lane arrow to Lane 3
    _wf_arrow(s, MX+1.1, 6.5, 8.68, 6.5, GREEN, 1.2)

    # ══ Lane 3: Generate & Stream ═══════════════════════════
    RX = 10.9   # center x of lane 3
    _wf_diamond(s, RX-1.0, 2.25, 2.0, 0.65, "Reasoning\nMode ON?", RGBColor(0x28,0x10,0x40), 9)
    _wf_label(s, "YES", RX-1.05, 2.48, PURPLE)
    _wf_arrow(s, RX-1.0, 2.58, RX-3.1, 2.58, PURPLE)
    _wf_box(s, RX-3.55, 2.33, 1.55, 0.5, "Reasoning\nModel (DeepSeek)", RGBColor(0x20,0x08,0x38), 9, PURPLE)
    _wf_arrow(s, RX-2.78, 2.83, RX-2.78, 3.3, PURPLE)
    _wf_box(s, RX-3.55, 3.3, 1.55, 0.5, "Analysis →\nFeed to Main", RGBColor(0x20,0x08,0x38), 9, PURPLE)

    _wf_label(s, "NO", RX+1.05, 2.48, ORANGE)
    _wf_arrow(s, RX+1.0, 2.58, RX+1.6, 2.58, ORANGE)
    # arrow down from diamond
    _wf_arrow(s, RX, 2.9, RX, 3.2, ORANGE)
    _wf_box(s, RX-1.0, 3.2, 2.0, 0.52, "Main LLM\n(Gemma / Qwen)", RGBColor(0x18,0x2A,0x18), 9.5, GREEN)
    _wf_arrow(s, RX, 3.72, RX, 4.0, ORANGE)
    _wf_diamond(s, RX-0.95, 4.0, 1.9, 0.6, "Main Model\nSucceeded?", RGBColor(0x2A,0x18,0x08), 9)
    _wf_label(s, "NO", RX+0.98, 4.2, PINK)
    _wf_arrow(s, RX+0.95, 4.3, RX+1.55, 4.3, PINK)
    _wf_box(s, RX+1.55, 4.05, 1.6, 0.5, "Fallback\nModel (Qwen)", RGBColor(0x2A,0x10,0x0A), 9, PINK)
    _wf_label(s, "YES", RX-1.1, 4.22, GREEN)
    _wf_arrow(s, RX, 4.6, RX, 4.9, GREEN)
    _wf_box(s, RX-1.0, 4.9, 2.0, 0.52, "Stream Tokens\nvia WebSocket", CARD2, 9.5)
    _wf_arrow(s, RX, 5.42, RX, 5.72, CYAN)
    _wf_box(s, RX-1.0, 5.72, 2.0, 0.52, "Update Conversation\nHistory + Memory", RGBColor(0x18,0x20,0x08), 9, YELLOW)
    _wf_arrow(s, RX, 6.24, RX, 6.54, CYAN)
    _wf_oval(s, RX-0.8, 6.54, 1.6, 0.45, "Response Delivered", GREEN, 9)

    # Legend
    add_text(s, "★  Dashed outlines = decision nodes   ●  Rounded = process step   ◉ = start/end",
             0.3, 7.22, 12.7, 0.35, size=9, color=MGRAY, align=PP_ALIGN.CENTER)


# ── Blueprint 2 : Document Ingestion & Query Workflow ─────────────────────────

def slide_wf_documents(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, GREEN, 0.07)

    add_text(s, "WORKFLOW  ②", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=GREEN)
    add_text(s, "Document Ingestion & Query Workflow", 0.3, 0.45, 12, 0.85, size=30, bold=True, color=WHITE)
    add_text(s, "From raw file upload to BM25-indexed retrieval and LLM-augmented response",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # ── TOP HALF: Ingestion pipeline (horizontal, left→right) ──────────────
    add_text(s, "INGESTION PATH", 0.3, 1.9, 4, 0.32, size=9.5, bold=True, color=GREEN)

    ing_steps = [
        (GREEN,   "File Upload\n(UI / CLI / Chat)"),
        (TEAL,    "Loader\nDetects Type"),
        (BLUE,    "Extract Text\n(PDF/DOCX/OCR/Audio)"),
        (CYAN,    "Chunk Text\n1000 ch / 200 overlap"),
        (PURPLE,  "SHA-256\nDedup Hash"),
        (ORANGE,  "Add to\nBM25 Index"),
        (YELLOW,  "Persist\npage_index.json"),
    ]
    box_w, box_h, gap = 1.65, 0.75, 0.15
    x = 0.3
    centers = []
    for col, lbl in ing_steps:
        _wf_box(s, x, 2.25, box_w, box_h, lbl, col, 8.5, BG if col in (YELLOW, GREEN, CYAN) else WHITE)
        cx = x + box_w / 2
        centers.append(cx)
        x += box_w + gap
    # arrows
    for i in range(len(centers) - 1):
        _wf_arrow(s, centers[i] + box_w/2 - 0.0, 2.625,
                  centers[i+1] - box_w/2 + 0.0, 2.625, GREEN, 1.2)
    # fix: arrows between boxes
    x = 0.3
    for i in range(len(ing_steps) - 1):
        _wf_arrow(s, x + box_w, 2.625, x + box_w + gap, 2.625, GREEN, 1.2)
        x += box_w + gap

    # Decision: duplicate?
    mid_x = centers[4]   # after hash step
    _wf_arrow(s, mid_x, 3.0, mid_x, 3.3, PURPLE)
    _wf_diamond(s, mid_x - 0.7, 3.3, 1.4, 0.52, "Duplicate\nChunk?", RGBColor(0x28,0x10,0x40), 8.5)
    _wf_label(s, "YES → SKIP", mid_x + 0.75, 3.5, PINK, 8)
    _wf_arrow(s, mid_x + 0.7, 3.56, mid_x + 1.3, 3.56, PINK)
    _wf_oval(s, mid_x + 1.3, 3.38, 1.4, 0.42, "Discard\nDuplicate", PINK, 8, BG)
    _wf_label(s, "NO → INDEX", mid_x - 0.75, 3.5, GREEN, 8)

    divider_line(s, 4.05)

    # ── BOTTOM HALF: Query path (horizontal) ───────────────────────────────
    add_text(s, "QUERY PATH", 0.3, 4.1, 4, 0.32, size=9.5, bold=True, color=CYAN)

    qry_steps = [
        (CYAN,    "User Query\nor File in Chat"),
        (BLUE,    "BM25 Search\n+ Term Expand"),
        (TEAL,    "Top-K Chunk\nRetrieval"),
        (GREEN,   "Context\nBuilder"),
        (ORANGE,  "Prompt\nAssembly"),
        (PURPLE,  "LLM\nGeneration"),
        (YELLOW,  "Streamed\nResponse"),
    ]
    x = 0.3
    for col, lbl in qry_steps:
        _wf_box(s, x, 4.45, box_w, box_h, lbl, col, 8.5, BG if col in (YELLOW, GREEN, CYAN) else WHITE)
        x += box_w + gap
    x = 0.3
    for i in range(len(qry_steps) - 1):
        _wf_arrow(s, x + box_w, 4.825, x + box_w + gap, 4.825, CYAN, 1.2)
        x += box_w + gap

    # Annotation boxes below
    annotations = [
        (0.3,   "Medical term\nexpansion built-in"),
        (2.1,   "top_k=3\n(configurable)"),
        (3.9,   "Preserves chunk\nboundary context"),
        (5.7,   "System prompt\n+ RAG + memory"),
        (7.5,   "Reasoning model\nfires if needed"),
        (9.3,   "Token streaming\nvia WebSocket"),
    ]
    for ax, note in annotations:
        add_text(s, note, ax, 5.32, 1.6, 0.55, size=8, color=MGRAY, align=PP_ALIGN.CENTER)

    # Supported formats box
    add_rounded_rect(s, 0.3, 6.05, 12.7, 1.22, fill=CARD, line=GREEN, line_w=1)
    add_text(s, "Supported File Types:", 0.5, 6.15, 2.8, 0.38, size=11, bold=True, color=GREEN)
    fmt_groups = [
        ("Documents", "PDF · DOCX · DOC · TXT · MD · Markdown"),
        ("Spreadsheets", "XLSX · XLS · CSV"),
        ("Images", "JPG · PNG · GIF · BMP · WebP · TIFF · SVG  (OCR optional)"),
        ("Audio", "MP3 · WAV · OGG · FLAC · M4A · AAC · WMA  (transcription optional)"),
    ]
    x = 0.5
    for cat, fmts in fmt_groups:
        add_text(s, cat + ":", x, 6.57, 1.4, 0.32, size=9.5, bold=True, color=LGRAY)
        add_text(s, fmts,      x, 6.9,  2.5, 0.32, size=9,   color=WHITE)
        x += 3.15


# ── Blueprint 3 : Analytics Workflow ─────────────────────────────────────────

def slide_wf_analytics(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, ORANGE, 0.07)

    add_text(s, "WORKFLOW  ③", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=ORANGE)
    add_text(s, "Data Analytics Workflow", 0.3, 0.45, 12, 0.85, size=32, bold=True, color=WHITE)
    add_text(s, "From raw data file to interactive charts, ML models, and predictions",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # Main vertical pipeline (center column)
    CX = 3.6
    steps = [
        (ORANGE, "Upload Data File\nCSV · Excel · JSON · TSV · PDF"),
        (BLUE,   "Data Loader\npandas.read_*() — all sheets"),
        (CYAN,   "Schema Detection\ndtype inference · null check"),
        (TEAL,   "Basic Statistics\nmean · median · std · quartiles"),
        (GREEN,  "Advanced Statistics\nskewness · kurtosis · IQR · correlations"),
        (PURPLE, "Outlier Detection\nIQR method  +  Z-score method"),
    ]
    y = 2.0
    step_centers = []
    for col, lbl in steps:
        _wf_box(s, CX - 1.35, y, 2.7, 0.62, lbl, col, 9, BG if col == ORANGE else WHITE)
        step_centers.append((CX, y + 0.31))
        if y < 5.5:
            _wf_arrow(s, CX, y + 0.62, CX, y + 0.77, col)
        y += 0.77

    # Decision: has numeric cols for ML?
    _wf_arrow(s, CX, y, CX, y + 0.2, PURPLE)
    _wf_diamond(s, CX - 0.9, y + 0.2, 1.8, 0.62, "Numeric Cols\nAvailable?", RGBColor(0x28,0x10,0x40), 9)
    _wf_label(s, "NO", CX - 1.0, y + 0.42, PINK)
    _wf_arrow(s, CX - 0.9, y + 0.51, CX - 1.8, y + 0.51, PINK)
    _wf_oval(s, CX - 3.2, y + 0.33, 1.35, 0.42, "Stats Only\nReport", PINK, 8.5, BG)
    _wf_label(s, "YES", CX + 1.0, y + 0.42, GREEN)
    _wf_arrow(s, CX + 0.9, y + 0.51, CX + 1.8, y + 0.51, GREEN)
    _wf_box(s, CX + 1.8, y + 0.28, 1.6, 0.5, "ML Model\nSelection", RGBColor(0x0A,0x28,0x18), 9, GREEN)

    # Right branch: ML models
    RX = 9.5
    add_text(s, "ML MODELS", RX - 0.5, 2.0, 3.0, 0.35, size=9.5, bold=True, color=PURPLE)
    ml_models = [
        (PURPLE, "Linear Regression\n(continuous target)"),
        (BLUE,   "Random Forest Regressor\n(non-linear)"),
        (CYAN,   "Logistic Regression\n(binary classification)"),
        (TEAL,   "Random Forest Classifier\n(multi-class)"),
        (GREEN,  "Time Series Trend\n(future forecasting)"),
    ]
    my = 2.4
    for col, lbl in ml_models:
        _wf_box(s, RX - 1.15, my, 2.3, 0.58, lbl, col, 8.5)
        if my < 5.5:
            _wf_arrow(s, RX, my + 0.58, RX, my + 0.68, col)
        my += 0.68

    # Converge to visualization
    _wf_arrow(s, RX, my, RX, my + 0.25, GREEN)
    _wf_box(s, RX - 1.15, my + 0.25, 2.3, 0.58, "Plotly Chart\nGeneration", RGBColor(0x18,0x28,0x08), 9, GREEN)
    # arrow from left pipeline converges too
    _wf_arrow(s, CX, 7.35, RX, 7.35, ORANGE)

    # Chart types row
    add_text(s, "VISUALIZATION OUTPUT", 0.3, 7.4, 5, 0.35, size=9.5, bold=True, color=YELLOW)
    charts = ["Histogram", "Scatter", "Line", "Bar", "Box", "Heatmap", "Pie", "Auto-Viz"]
    x = 0.3
    for ct in charts:
        add_rounded_rect(s, x, 7.78, 1.45, 0.38, fill=CARD2, line=ORANGE, line_w=1)
        add_text(s, ct, x + 0.05, 7.82, 1.35, 0.3, size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        x += 1.58


# ── Blueprint 4 : MCP Tool-Call Workflow ─────────────────────────────────────

def slide_wf_mcp(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, PURPLE, 0.07)

    add_text(s, "WORKFLOW  ④", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=PURPLE)
    add_text(s, "MCP Tool-Call Workflow", 0.3, 0.45, 12, 0.85, size=32, bold=True, color=WHITE)
    add_text(s, "How RAGenie connects to external MCP servers and executes tool calls",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # ── Left: Connection Setup (one-time) ──────────────────
    add_text(s, "SETUP  (one-time per server)", 0.3, 1.92, 5.5, 0.35, size=9.5, bold=True, color=CYAN)
    setup = [
        (CYAN,   "User Registers MCP Server\n(UI: name · transport · URL/command)"),
        (BLUE,   "Config Saved\nto SQLite DB"),
        (TEAL,   "MCPClientManager\ninitialises connection"),
        (GREEN,  "Handshake: list_tools()\nFetch all tool schemas"),
        (YELLOW, "Register tools in LLM Agent\nas  server/tool_name"),
    ]
    y = 2.35
    for col, lbl in setup:
        _wf_box(s, 0.3, y, 3.5, 0.58, lbl, col, 8.5, BG if col in (CYAN, YELLOW, GREEN) else WHITE)
        if y < 4.9:
            _wf_arrow(s, 2.05, y + 0.58, 2.05, y + 0.73, col)
        y += 0.73

    # vertical divider
    add_rect(s, 4.15, 1.9, 0.04, 5.5, fill=MGRAY)

    # ── Right: Per-message tool-call flow ──────────────────
    add_text(s, "PER MESSAGE  (runtime)", 4.3, 1.92, 8.5, 0.35, size=9.5, bold=True, color=ORANGE)

    RX = 8.0
    _wf_oval(s, RX - 0.9, 2.35, 1.8, 0.48, "User Message\nReceived", ORANGE, 9, BG)
    _wf_arrow(s, RX, 2.83, RX, 3.1, ORANGE)
    _wf_box(s, RX - 1.1, 3.1, 2.2, 0.55, "LLM Agent\nProcesses Prompt", BLUE, 9)
    _wf_arrow(s, RX, 3.65, RX, 3.92, ORANGE)
    _wf_diamond(s, RX - 1.0, 3.92, 2.0, 0.62, "Needs External\nTool?", RGBColor(0x28,0x18,0x08), 9)
    _wf_label(s, "NO", RX - 1.1, 4.14, LGRAY)
    _wf_arrow(s, RX - 1.0, 4.23, RX - 1.9, 4.23, LGRAY)
    _wf_box(s, RX - 3.7, 3.98, 1.75, 0.5, "Direct LLM\nResponse", CARD2, 8.5)
    _wf_label(s, "YES", RX + 1.1, 4.14, CYAN)
    _wf_arrow(s, RX + 1.0, 4.23, RX + 1.8, 4.23, CYAN)
    _wf_box(s, RX + 1.8, 3.98, 1.8, 0.5, "Select Tool:\nserver/tool_name", RGBColor(0x08,0x18,0x38), 8.5, CYAN)

    _wf_arrow(s, RX, 4.54, RX, 4.82, ORANGE)
    _wf_diamond(s, RX - 1.0, 4.82, 2.0, 0.62, "Transport\nType?", RGBColor(0x18,0x10,0x38), 9)

    # Three transport branches
    for dx, label, col, transport in [
        (-2.5, "stdio", TEAL,   "Launch subprocess\nstdin/stdout IPC"),
        (0,    "SSE",   BLUE,   "HTTP GET\nServer-Sent Events"),
        (+2.5, "HTTP",  PURPLE, "HTTP POST\nStreamable HTTP"),
    ]:
        tx = RX + dx
        _wf_label(s, label, tx, 5.34, col, 8.5)
        _wf_arrow(s, tx, 5.44, tx, 5.68, col)
        _wf_box(s, tx - 0.85, 5.68, 1.7, 0.52, transport, RGBColor(0x0A,0x18,0x30), 8.5, col)
        _wf_arrow(s, tx, 6.2, tx, 6.42, col)

    # Converge
    for dx in [-2.5, 0, 2.5]:
        _wf_arrow(s, RX + dx, 6.42, RX, 6.42, ORANGE)
    _wf_box(s, RX - 1.1, 6.42, 2.2, 0.55, "Return Tool Result\nto LLM Agent", CARD2, 9)
    _wf_arrow(s, RX, 6.97, RX, 7.2, ORANGE)
    _wf_oval(s, RX - 0.9, 7.2, 1.8, 0.42, "Final Response\nStreamed to User", GREEN, 9, BG)

    # UUID registry note
    add_text(s, "Registry: UUID-keyed server store  ·  LLM sees: server_name/tool_name  ·  Tools rebuild live on connect/disconnect",
             0.3, 7.3, 7.6, 0.35, size=8.5, color=MGRAY)


# ── Blueprint 5 : Auth & Security Workflow ───────────────────────────────────

def slide_wf_auth(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, PINK, 0.07)

    add_text(s, "WORKFLOW  ⑤", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=PINK)
    add_text(s, "Authentication & Request Security Workflow", 0.3, 0.45, 12, 0.85, size=28, bold=True, color=WHITE)
    add_text(s, "JWT token lifecycle, rate limiting, and request validation pipeline",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # ── Left: Token lifecycle ──────────────────────────────
    add_text(s, "TOKEN LIFECYCLE", 0.3, 1.92, 5, 0.35, size=9.5, bold=True, color=PINK)

    CX = 2.1
    _wf_oval(s, CX - 0.9, 2.35, 1.8, 0.48, "User: POST\n/auth/login", PINK, 9, BG)
    _wf_arrow(s, CX, 2.83, CX, 3.1, PINK)
    _wf_box(s, CX - 1.05, 3.1, 2.1, 0.55, "Verify credentials\nbcrypt hash check", RGBColor(0x28,0x08,0x18), 9, PINK)
    _wf_arrow(s, CX, 3.65, CX, 3.92, PINK)
    _wf_diamond(s, CX - 0.9, 3.92, 1.8, 0.62, "Credentials\nValid?", RGBColor(0x28,0x08,0x18), 9)
    _wf_label(s, "NO", CX - 1.0, 4.14, ORANGE)
    _wf_arrow(s, CX - 0.9, 4.23, CX - 1.65, 4.23, ORANGE)
    _wf_oval(s, CX - 2.8, 4.05, 1.1, 0.42, "401\nUnauthorized", ORANGE, 8.5, BG)
    _wf_label(s, "YES", CX + 1.0, 4.14, GREEN)
    _wf_arrow(s, CX, 4.54, CX, 4.82, GREEN)
    _wf_box(s, CX - 1.05, 4.82, 2.1, 0.6,
            "Issue:\n• Access Token (30 min)\n• Refresh Token (7 days)", RGBColor(0x0A,0x28,0x18), 9, GREEN)
    _wf_arrow(s, CX, 5.42, CX, 5.7, GREEN)
    _wf_box(s, CX - 1.05, 5.7, 2.1, 0.55, "Client stores tokens\nLocally (localStorage)", CARD2, 9)
    _wf_arrow(s, CX, 6.25, CX, 6.52, GREEN)
    _wf_diamond(s, CX - 0.9, 6.52, 1.8, 0.62, "Access Token\nExpired?", RGBColor(0x28,0x18,0x08), 9)
    _wf_label(s, "YES", CX + 1.0, 6.74, YELLOW)
    _wf_arrow(s, CX + 0.9, 6.83, CX + 1.6, 6.83, YELLOW)
    _wf_box(s, CX + 1.6, 6.6, 1.8, 0.5, "POST /auth/refresh\n→ new Access Token", RGBColor(0x28,0x20,0x08), 8.5, YELLOW)
    _wf_label(s, "NO", CX - 1.0, 6.74, GREEN)
    _wf_arrow(s, CX - 0.9, 6.83, CX - 1.6, 6.83, GREEN)
    _wf_oval(s, CX - 2.85, 6.65, 1.2, 0.42, "Request\nProceeds", GREEN, 8.5, BG)

    # Divider
    add_rect(s, 4.35, 1.9, 0.04, 5.5, fill=MGRAY)

    # ── Right: Request pipeline ────────────────────────────
    add_text(s, "REQUEST PIPELINE  (every API call)", 4.5, 1.92, 8.5, 0.35, size=9.5, bold=True, color=CYAN)

    stages = [
        (CYAN,   "Incoming HTTP / WS Request"),
        (BLUE,   "Security Headers Applied\n(CSP · HSTS · X-Frame)"),
        (TEAL,   "CORS Origin Check\n(whitelist validation)"),
        (ORANGE, "Rate Limiter (slowapi)\n60 req/min · 10 uploads/hr · 30 WS/min"),
        (PURPLE, "JWT Verification\n(when auth.enabled=true)"),
        (GREEN,  "Request Size Check\n(max 30 MB · 10 000 char WS)"),
        (YELLOW, "Content Validation\nFile type · Schema"),
        (PINK,   "Audit Log Written\nlogs/audit.log  (redacted)"),
    ]
    rx = 6.8
    y = 2.35
    for col, lbl in stages:
        _wf_box(s, rx - 0.0, y, 6.1, 0.54, lbl, col, 9, BG if col in (CYAN, YELLOW, GREEN) else WHITE)
        if y < 6.6:
            _wf_arrow(s, rx + 3.05, y + 0.54, rx + 3.05, y + 0.64, col)
        y += 0.64

    _wf_arrow(s, rx + 3.05, y, rx + 3.05, y + 0.2, GREEN)
    _wf_oval(s, rx + 1.55, y + 0.2, 3.0, 0.45, "Handler Executes  →  Response", GREEN, 9.5, BG)


# ── Blueprint 6 : Memory & Proactive Workflow ────────────────────────────────

def slide_wf_memory(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    accent_bar(s, YELLOW, 0.07)

    add_text(s, "WORKFLOW  ⑥", 0.3, 0.1, 4, 0.4, size=11, bold=True, color=YELLOW)
    add_text(s, "Memory, Learning & Proactive Workflow", 0.3, 0.45, 12, 0.85, size=30, bold=True, color=WHITE)
    add_text(s, "How RAGenie remembers, learns from feedback, and proactively assists users",
             0.3, 1.3, 12, 0.45, size=14, color=LGRAY)
    divider_line(s, 1.8)

    # ── Left: Memory write path ────────────────────────────
    add_text(s, "MEMORY WRITE", 0.3, 1.92, 4, 0.35, size=9.5, bold=True, color=YELLOW)
    mem_steps = [
        (YELLOW, "Conversation\nTurn Completes"),
        (ORANGE, "Extract Important\nContext Phrases"),
        (TEAL,   "Embed in\nMemory Store (SQLite)"),
        (GREEN,  "Prune to\nmax_context_items=8"),
    ]
    y = 2.35
    for col, lbl in mem_steps:
        _wf_box(s, 0.3, y, 2.8, 0.58, lbl, col, 9, BG if col == YELLOW else WHITE)
        if y < 4.5:
            _wf_arrow(s, 1.7, y + 0.58, 1.7, y + 0.72, col)
        y += 0.72

    # ── Center: Memory read path ──────────────────────────
    add_text(s, "MEMORY READ", 3.4, 1.92, 4, 0.35, size=9.5, bold=True, color=CYAN)
    read_steps = [
        (CYAN,   "New User\nMessage"),
        (BLUE,   "Query Memory\nStore (relevance)"),
        (PURPLE, "Inject Top-8\nMemory Items"),
        (TEAL,   "Build Prompt\nwith Memory Context"),
    ]
    y = 2.35
    for col, lbl in read_steps:
        _wf_box(s, 3.4, y, 2.8, 0.58, lbl, col, 9, WHITE)
        if y < 4.5:
            _wf_arrow(s, 4.8, y + 0.58, 4.8, y + 0.72, col)
        y += 0.72

    # ── Right: Learning feedback ──────────────────────────
    add_text(s, "LEARNING FEEDBACK", 6.5, 1.92, 4, 0.35, size=9.5, bold=True, color=GREEN)
    _wf_oval(s, 6.5, 2.35, 2.8, 0.48, "AI Response\nDisplayed", GREEN, 9, BG)
    _wf_arrow(s, 7.9, 2.83, 7.9, 3.1, GREEN)
    _wf_diamond(s, 6.9, 3.1, 2.0, 0.62, "User Gives\nFeedback?", RGBColor(0x08,0x28,0x18), 9)
    _wf_label(s, "NO", 6.78, 3.32, LGRAY)
    _wf_arrow(s, 6.9, 3.41, 6.1, 3.41, LGRAY)
    _wf_oval(s, 5.1, 3.22, 0.95, 0.42, "No\nChange", MGRAY, 8.5, WHITE)
    _wf_label(s, "YES", 9.0, 3.32, GREEN)
    _wf_arrow(s, 8.9, 3.41, 9.6, 3.41, GREEN)
    _wf_diamond(s, 9.6, 3.22, 1.85, 0.58, "👍 or\n👎?", RGBColor(0x08,0x28,0x18), 9)
    _wf_label(s, "👍 +0.10", 9.45, 3.7, GREEN, 8)
    _wf_arrow(s, 9.45, 3.8, 8.8, 3.97, GREEN)
    _wf_box(s, 7.5, 3.97, 2.8, 0.5, "Update Retrieval\nScore (SQLite)", RGBColor(0x0A,0x20,0x0A), 9, GREEN)
    _wf_label(s, "👎 -0.08", 11.65, 3.7, PINK, 8)
    _wf_arrow(s, 11.55, 3.8, 10.5, 3.97, PINK)
    _wf_arrow(s, 7.9, 4.47, 7.9, 4.72, GREEN)
    _wf_box(s, 6.5, 4.72, 2.8, 0.55, "Scores influence\nfuture BM25 ranking", CARD2, 9)

    # ── Proactive system ──────────────────────────────────
    divider_line(s, 5.5)
    add_text(s, "PROACTIVE CAPABILITIES", 0.3, 5.56, 6, 0.35, size=9.5, bold=True, color=TEAL)

    proactive_flow = [
        (TEAL,   "Background\nScheduler"),
        (BLUE,   "Check Time\nvs quiet hours"),
        (CYAN,   "Fetch News\n+ Memory"),
        (GREEN,  "Generate\nDaily Briefing"),
        (YELLOW, "Push Nudge\nto User"),
    ]
    x = 0.3
    for col, lbl in proactive_flow:
        _wf_box(s, x, 5.97, 2.3, 0.62, lbl, col, 8.5, BG if col in (TEAL, YELLOW) else WHITE)
        if x < 9.5:
            _wf_arrow(s, x + 2.3, 6.28, x + 2.45, 6.28, col)
        x += 2.45

    # Config note
    add_text(s, "Configurable:  briefing_hour=9  ·  cycle_interval=30 min  ·  quiet_hours 22:00–08:00  ·  enabled: true/false",
             0.3, 6.72, 12.7, 0.38, size=9.5, color=MGRAY, align=PP_ALIGN.CENTER)

    add_text(s, "News aggregator feeds the proactive briefing — articles are optionally ingested into the RAG index too.",
             0.3, 7.1, 12.7, 0.35, size=9, color=MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation(output_path="RAGenie_Overview.pptx"):
    prs = new_prs()

    print("Building slides...")
    slide_cover(prs);           print("  ✓ Cover")
    slide_what_is(prs);         print("  ✓ What is RAGenie?")
    slide_architecture(prs);    print("  ✓ System Architecture")
    slide_rag_flow(prs);        print("  ✓ RAG Pipeline")
    slide_multi_model(prs);     print("  ✓ Multi-Model Architecture")
    slide_capabilities(prs);    print("  ✓ Core Capabilities")
    slide_analytics(prs);       print("  ✓ Data Analytics")
    slide_mcp(prs);             print("  ✓ MCP Integration")
    slide_security(prs);        print("  ✓ Security & Auth")
    slide_memory_learning(prs); print("  ✓ Memory & Learning")
    slide_news(prs);            print("  ✓ News Aggregator")
    slide_tech_stack(prs);      print("  ✓ Technology Stack")
    slide_api(prs);             print("  ✓ API Reference")
    slide_getting_started(prs); print("  ✓ Getting Started")
    slide_requirements(prs);    print("  ✓ System Requirements")
    slide_differentiators(prs); print("  ✓ Why RAGenie?")

    # ── Blueprint / Workflow slides ──
    slide_wf_chat(prs);       print("  ✓ [Blueprint] Chat Lifecycle")
    slide_wf_documents(prs);  print("  ✓ [Blueprint] Document Workflow")
    slide_wf_analytics(prs);  print("  ✓ [Blueprint] Analytics Workflow")
    slide_wf_mcp(prs);        print("  ✓ [Blueprint] MCP Tool-Call Workflow")
    slide_wf_auth(prs);       print("  ✓ [Blueprint] Auth & Security Workflow")
    slide_wf_memory(prs);     print("  ✓ [Blueprint] Memory & Proactive Workflow")

    slide_closing(prs);         print("  ✓ Closing")

    prs.save(output_path)
    print(f"\n✅  Saved: {output_path}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build_presentation()
